"""
This file contains two functions which make a powerpoint of the visualisations in the tool.
"""
from datetime import date, datetime
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_LEGEND_POSITION
import numpy as np


def make_title(slide, case, target, title_page=False):
    """
    This function adds a title to a slide
    :param slide: selected slide where the title is needed
    :param case: selected case
    :param target: the desired slide where the title is wanted
    :param title_page: indicate if title is made on title page
    """
    title = slide.shapes.title
    if title_page:
        title.left, title.top, title.width, title.height = Cm(5.9), Cm(5.7), Cm(21.6), Cm(4.1)
    else:
        title.left, title.top, title.width, title.height = Cm(5.1), Cm(0.7), Cm(22.9), Cm(3.2)
    try:
        title.text = case.input_dict["text_element_value"][
            list(case.input_dict["text_elements"]).index("title_" + target)
        ]
    except:
        title.text = "Not defined in template"
        title.text.paragraphs[0].font.italic = True
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)


def make_subtitle(slide):
    """
    This function adds a subtitle to a slide
    :param slide: selected slide where the title is needed
    """
    date_today = str(date.today().strftime("%d %B %Y"))
    subtitle = slide.placeholders[1]
    subtitle.top, subtitle.left, subtitle.width, subtitle.height = Cm(10.8), Cm(7.6), Cm(17.8), Cm(4.9)
    subtitle.text = f"Responsible business decision making \n {date_today}"


def make_introduction(slide, case, target):
    """
    This function adds an introduction to a slide
    :param slide: selected slide where the title is needed
    :param case: selected case
    :param target: the desired slide where the title is wanted
    """
    tx_box = slide.shapes.add_textbox(Cm(5.1), Cm(3.3), Cm(22.8), Cm(2.6))
    text_frame = tx_box.text_frame
    text_frame.word_wrap = True
    try:
        text_frame.text = case.input_dict["text_element_value"][
            list(case.input_dict["text_elements"]).index("intro_" + target)
        ]
    except:
        text_frame.text = "Not defined in template"
        text_frame.paragraphs[0].font.italic = True


def make_strategic_challenge(slide, case):
    """
    This function adds an introduction to a slide
    :param slide: selected slide where the title is needed
    :param case: the selected case
    """
    input_info = case.input_dict
    tx_box = slide.shapes.add_textbox(Cm(5.1), Cm(4.6), Cm(22.8), Cm(5))
    try:
        tx_box.text_frame.text = input_info["text_element_value"][
            list(case.input_dict["text_elements"]).index("strategic_challenge")
        ]
    except:
        tx_box.text_frame.text = "Not defined in template"
        tx_box.text_frame.paragraphs[0].font.italic = True
    tx_box.text_frame.paragraphs[0].font.size = Pt(28)
    tx_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def iter_cells(table):
    """
    This function makes it possible to iterate over all cells in a table
    :param table: selected table
    """
    for row in table.rows:
        for cell in row.cells:
            yield cell


def change_format_cells(table, format_size):
    """
    This function makes it possible to change the format of all cells in a table
    :param table: selected table
    :param format_size: the desired format size of the cells in a table
    """
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(format_size)


def add_page_number(slide, page_number):
    """
    This function adds a page number to a slide
    :param slide: selected slide where the title is needed
    :param page_number: the number on the page
    """
    x_loc_page_number, y_loc_page_number, x_weight_page_number, y_weight_page_number = (
        Cm(31.9),
        Cm(17.7),
        Cm(0.8),
        Cm(1),
    )
    tx_box = slide.shapes.add_textbox(x_loc_page_number, y_loc_page_number, x_weight_page_number, y_weight_page_number)
    tx_box.text_frame.text = page_number
    tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(120, 120, 120)


def make_slides(case, scenario):
    """
    This function present different visualisation on different slides in a powerpoint format
    :param case: selected case with all start and calculated values
    :param scenario: the desired scenario which is used in the visualisations
    """
    prs = Presentation()
    # set width to full width.
    prs.slide_width = Inches(13)
    output_info = case.output_dict[scenario]
    input_info = case.input_dict
    # Formatting position
    x_loc_graph, y_loc_graph, x_weight_graph, y_weight_graph = Cm(3.62), Cm(3.9), Cm(25.5), Cm(14.8)
    x_loc_table, y_loc_table, x_weight_table, y_weight_table = Cm(3.62), Cm(6.2), Cm(25.5), Cm(0.5)

    # Generate title slide
    slide_title = prs.slides.add_slide(prs.slide_layouts[0])
    make_title(slide_title, case, "front_page", True)
    make_subtitle(slide_title)

    # Generate strategic challenge slide
    slide_challenge = prs.slides.add_slide(prs.slide_layouts[5])
    make_title(slide_challenge, case, "strategic_challenge")
    make_strategic_challenge(slide_challenge, case)
    add_page_number(slide_challenge, "1")

    # Create the table chart slide key outputs
    slide_key_outputs = prs.slides.add_slide(prs.slide_layouts[5])
    make_title(slide_key_outputs, case, "key_outputs")
    make_introduction(slide_key_outputs, case, "key_outputs")
    add_page_number(slide_key_outputs, "2")
    #  Data for table
    input_info_ko = input_info["key_outputs"]
    input_info_ko_th = input_info["key_output_theme"]
    shape = slide_key_outputs.shapes.add_table(
        (len(input_info_ko) + 1), 2, x_loc_table, y_loc_table, x_weight_table, y_weight_table
    )
    table = shape.table
    headers_sum = ["Key outputs", "Theme"]
    for column, item in enumerate(headers_sum):
        cell = table.cell(0, column)
        cell.text = headers_sum[column]
    for column, table_name in enumerate([input_info_ko, input_info_ko_th]):
        for row in range(1, (len(input_info_ko) + 1)):
            cell = table.cell(row, column)
            cell.text = str(table_name[row - 1])
    change_format_cells(table, 12)

    # Create the table chart slide decision maker's options
    slide_dmo = prs.slides.add_slide(prs.slide_layouts[5])
    make_title(slide_dmo, case, "dmo")
    make_introduction(slide_dmo, case, "dmo")
    add_page_number(slide_dmo, "3")
    #  Data for table
    input_info_dmo = input_info["decision_makers_options"]
    input_info_dmo_values = input_info["decision_makers_option_value"]
    input_info_internal_value = input_info["internal_variable_inputs"]
    shape = slide_dmo.shapes.add_table(
        (len(input_info_internal_value) + 1),
        len(input_info_dmo) + 1,
        x_loc_table,
        y_loc_table,
        x_weight_table,
        y_weight_table,
    )
    table = shape.table
    headers_sum = np.append(["Internal variable input"], input_info_dmo)
    #  Fill headers
    for column, item in enumerate(headers_sum):
        cell = table.cell(0, column)
        cell.text = headers_sum[column]
    #  Fill left column
    for row in range(1, len(input_info_internal_value) + 1):
        cell = table.cell(row, 0)
        cell.text = input_info_internal_value[row - 1]
    #  Fill rest of the table
    for column in range(1, len(input_info_dmo) + 1):
        for row in range(1, (len(input_info_internal_value) + 1)):
            cell = table.cell(row, column)
            value = float(input_info_dmo_values[column - 1][row - 1])
            cell.text = f"{value:.2f}"
    change_format_cells(table, 12)

    # Create the table chart slide the scenarios
    slide_scenarios = prs.slides.add_slide(prs.slide_layouts[5])
    make_title(slide_scenarios, case, "scenarios")
    make_introduction(slide_scenarios, case, "scenarios")
    add_page_number(slide_scenarios, "4")
    #  Data for table
    input_info_scen = input_info["scenarios"]
    input_info_scen_values = input_info["scenario_value"]
    input_info_external_value = input_info["external_variable_inputs"]
    shape = slide_scenarios.shapes.add_table(
        (len(input_info_external_value) + 1),
        len(input_info_scen) + 1,
        x_loc_table,
        y_loc_table,
        x_weight_table,
        y_weight_table,
    )
    table = shape.table
    headers_sum = np.append(["External variable input"], input_info_scen)
    #  Fill headers
    for column_header, item in enumerate(headers_sum):
        cell = table.cell(0, column_header)
        cell.text = headers_sum[column_header]
    #  Fill left column
    for row in range(1, len(input_info_external_value) + 1):
        cell = table.cell(row, 0)
        cell.text = input_info_external_value[row - 1]
    #  Fill rest of the table
    for column in range(1, len(input_info_scen) + 1):
        for row in range(1, (len(input_info_external_value) + 1)):
            cell = table.cell(row, column)
            value = float(input_info_scen_values[column - 1][row - 1])
            cell.text = f"{value:.2f}"
    change_format_cells(table, 12)

    # Create the table chart slide fixed inputs
    slide_fixed_inputs = prs.slides.add_slide(prs.slide_layouts[5])
    make_title(slide_fixed_inputs, case, "fixed_inputs")
    # make_introduction(slide_fixed_inputs, case, 'fixed_inputs')
    add_page_number(slide_fixed_inputs, "5")
    #  Data for table
    input_info_fi = input_info["fixed_inputs"]
    input_info_fi_val = input_info["fixed_input_value"]
    shape = slide_fixed_inputs.shapes.add_table(
        (len(input_info_fi) + 1), 2, x_loc_table, y_loc_table, x_weight_table, y_weight_table
    )
    table = shape.table
    headers_sum = ["Fixed Inputs", "Value"]
    for column, item in enumerate(headers_sum):
        cell = table.cell(0, column)
        cell.text = headers_sum[column]
    for column, table_name in enumerate([input_info_fi, input_info_fi_val]):
        for row in range(1, (len(input_info_fi) + 1)):
            cell = table.cell(row, column)
            cell.text = str(table_name[row - 1])
    change_format_cells(table, 12)

    # Generate dependency graph slide
    slide_dep_graph = prs.slides.add_slide(prs.slide_layouts[5])
    make_title(slide_dep_graph, case, "dependency_graph")
    slide_dep_graph.shapes.add_picture("MicrosoftTeams-image.png", x_loc_table, y_loc_table, x_weight_table, Cm(12.5))

    add_page_number(slide_dep_graph, "6")

    # Generate weighted appreciations slide
    slide_weighted_app = prs.slides.add_slide(prs.slide_layouts[5])
    make_title(slide_weighted_app, case, "weighted_graph")
    add_page_number(slide_weighted_app, "7")
    #  Fill chart with data
    chart_data = CategoryChartData()
    chart_data.categories = output_info.keys()
    for series in list(output_info[list(output_info.keys())[0]]["weighted_appreciations"].keys()):
        values = [output_info[category]["weighted_appreciations"][series] for category in list(output_info.keys())]
        chart_data.add_series(series, values)
    #  implement chart on slide
    chart = slide_weighted_app.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, x_loc_graph, y_loc_graph, x_weight_graph, y_weight_graph, chart_data
    ).chart
    #  color series on chart
    colors = [
        "#D04A02",
        "#EB8C00",
        "#FFB600",
        "#295477",
        "#299D8F",
        "#D04A02",
        "#EB8C00",
        "#FFB600",
        "#295477",
        "#299D8F",
        "#D04A02",
        "#EB8C00",
        "#FFB600",
        "#295477",
        "#299D8F",
    ]
    for index, series in enumerate(chart.series):
        fill = series.format.fill
        fill.solid()
        color_used = colors[index].lstrip("#")
        rgb = tuple(int(color_used[i : i + 2], 16) for i in (0, 2, 4))
        fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
    #  layout chart change
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(9)
    value_axis = chart.value_axis
    value_axis.maximum_scale = 100
    value_axis.major_unit = 20
    #  titles
    chart.chart_title.text_frame.text = "Values of weighted appreciations"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    category_axis_title = chart.value_axis
    category_axis_title.axis_title.text_frame.text = "Appreciation value"
    category_axis_title.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
    #  for implementing legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.horz_offset = 0.0
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
    return prs


def create_report(case, scenario, path):
    """
    This function saves the created report out of the function make_powerpoint function in a desired location
    :param case: selected case with all start and calculated values
    :param scenario: the desired scenario which is used in the visualisations
    :param path: the desired location where the report is saved
    """
    prs = make_slides(
        case,
        scenario,
    )
    date_year = str(datetime.now().strftime("%Y-%m-%d"))
    date_hour = str(datetime.now().strftime("%Hh%Mm%Ss"))
    filename = "tRBS_" + case.name + f'_{date_year + "_" + date_hour}.pptx'
    prs.save(str(path) + "/" + filename)
    text_finished = "The powerpoint is generated and located at " + str(path) + "/" + filename
    return text_finished
