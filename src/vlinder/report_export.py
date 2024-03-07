"""
This file contains two functions which make a report of the information performed in the tool.
"""
import os
from pathlib import Path
from datetime import date, datetime
from random import random
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_LEGEND_POSITION
import numpy as np


def make_subtitle(slide) -> None:
    """
    This function adds a subtitle to a slide
    :param slide: selected slide where the title is needed
    """
    date_today = str(date.today().strftime("%d %B %Y"))
    subtitle = slide.placeholders[1]
    subtitle.top, subtitle.left, subtitle.width, subtitle.height = Cm(10.8), Cm(7.6), Cm(17.8), Cm(4.9)
    subtitle.text = f"Responsible business decision making \n {date_today}"


def iter_cells(table) -> None:
    """
    This function makes it possible to iterate over all cells in a table
    :param table: selected table
    """
    for row in table.rows:
        for cell in row.cells:
            yield cell


def change_format_cells(table, format_size) -> None:
    """
    This function makes it possible to change the format of all cells in a table
    :param table: selected table
    :param format_size: the desired format size of the cells in a table
    """
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(format_size)


class ReportExporter:
    """
    This class deals with the transformation into a different format and export of output of an RBS case.
    """

    def __init__(self, output_path, name, input_dict, output_dict):
        self.output_path = Path(output_path)
        self.folder_name = ""
        self.name = name
        self.input_dict = input_dict
        self.output_dict = output_dict
        self.page_number = 1

    def make_title(self, slide, target, title_page=False, scenario="", pos_series="", key_output="") -> None:
        """
        This function adds a title to a slide
        :param slide: selected slide where the title is needed
        :param target: the desired slide where the title is wanted
        :param title_page: indicate if title is made on title page
        :param scenario: scenario of the output dict
        :param pos_series: position in series
        :param key_output: the selected key_output
        """
        title = slide.shapes.title
        if title_page:
            title.left, title.top, title.width, title.height = Cm(5.9), Cm(5.7), Cm(21.6), Cm(4.1)
        else:
            title.left, title.top, title.width, title.height = Cm(5.1), Cm(0.7), Cm(22.9), Cm(3.2)
        if "title_" + target in self.input_dict["case_text_elements"]:
            text_element = "case_text_element"
        else:
            text_element = "generic_text_element"
        if "title_" + target in self.input_dict[text_element + "s"]:
            if str(
                    self.input_dict[text_element + "_value"][
                        list(self.input_dict[text_element + "s"]).index("title_" + target)
                    ]
            ) != str("nan"):
                title.text = (
                        self.input_dict[text_element + "_value"][
                            list(self.input_dict[text_element + "s"]).index("title_" + target)
                        ]
                        + scenario
                        + pos_series
                        + key_output
                )
            else:
                title.text = "Not defined in template"
                title.text_frame.paragraphs[0].runs[0].font.italic = True
                title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 0, 0)
        else:
            title.text = "Not defined in template"
            title.text_frame.paragraphs[0].runs[0].font.italic = True
            title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 0, 0)
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)

    def make_introduction(self, slide, target) -> None:
        """
        This function adds an introduction to a slide
        :param slide: selected slide where the title is needed
        :param target: the desired slide where the title is wanted
        """
        tx_box = slide.shapes.add_textbox(Cm(5.1), Cm(3.3), Cm(22.8), Cm(2.6))
        text_frame = tx_box.text_frame
        text_frame.word_wrap = True
        if "intro_" + target in self.input_dict["case_text_elements"]:
            text_element = "case_text_element"
        else:
            text_element = "generic_text_element"
        if "intro_" + target in self.input_dict[text_element + "s"]:
            if str(
                    self.input_dict[text_element + "_value"][
                        list(self.input_dict[text_element + "s"]).index("intro_" + target)
                    ]
            ) != str("nan"):
                text_frame.text = self.input_dict[text_element + "_value"][
                    list(self.input_dict[text_element + "s"]).index("intro_" + target)
                ]
            else:
                text_frame.text = "Not defined in template"
                text_frame.paragraphs[0].runs[0].font.italic = True
                text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 0, 0)
        else:
            text_frame.text = "Not defined in template"
            text_frame.paragraphs[0].runs[0].font.italic = True
            text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 0, 0)

    def make_strategic_challenge(self, slide) -> None:
        """
        This function adds an introduction to a slide
        :param slide: selected slide where the title is needed
        """
        tx_box = slide.shapes.add_textbox(Cm(5.1), Cm(4.6), Cm(22.8), Cm(5))
        tx_box.text_frame.word_wrap = True
        if "strategic_challenge" in self.input_dict["case_text_elements"]:
            if str(
                    self.input_dict["case_text_element_value"][
                        list(self.input_dict["case_text_elements"]).index("strategic_challenge")
                    ]
            ) != str("nan"):
                tx_box.text = self.input_dict["case_text_element_value"][
                    list(self.input_dict["case_text_elements"]).index("strategic_challenge")
                ]
            else:
                tx_box.text = "Not defined in template"
                tx_box.text_frame.paragraphs[0].runs[0].font.italic = True
                tx_box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 0, 0)
        else:
            tx_box.text = "Not defined in template"
            tx_box.text_frame.paragraphs[0].runs[0].font.italic = True
            tx_box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 0, 0)
        tx_box.text_frame.paragraphs[0].font.size = Pt(28)
        tx_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def make_table_2_col(
            self, slide, col_left, col_right, col_left_header, col_right_header, first_element=0, last_element=10
    ) -> None:
        """
        This function makes it possible to iterate over all cells in a table
        :param slide: selected slide where the title is needed
        :param col_left: left column of a table
        :param col_right: right column of a table
        :param col_left_header: header of left column of a table
        :param col_right_header: header of right column of a table
        :param first_element: position of the first element in table
        :param last_element: position of the last element in table
        """
        input_info_col_left = self.input_dict[col_left][first_element:last_element]
        input_info_col_right = self.input_dict[col_right][first_element:last_element]
        shape = slide.shapes.add_table((len(input_info_col_left) + 1), 2, Cm(3.62), Cm(6.2), Cm(25.5), Cm(0.5))
        table = shape.table
        headers_total = [col_left_header, col_right_header]
        for column, item in enumerate(headers_total):
            cell = table.cell(0, column)
            cell.text = headers_total[column]
        for column, table_name in enumerate([input_info_col_left, input_info_col_right]):
            for row in range(1, (len(input_info_col_left) + 1)):
                cell = table.cell(row, column)
                value = table_name[row - 1]
                if isinstance(value, float):
                    cell.text = '{:,.2f}'.format(value).replace(',', '*').replace('.', ',').replace('*', '.')
                else:
                    cell.text = value
        change_format_cells(table, 12)

    def make_table_n_col(self, slide, col_names, col_values, row_names, left_col_header) -> None:
        """
        This function makes it possible to iterate over all cells in a table
        :param slide: selected slide where the title is needed
        :param col_names: column names of a table
        :param col_values: values in the columns
        :param row_names: names of all the values in the left column
        :param left_col_header: header of left column of a table
        """
        input_info_col_names = self.input_dict[col_names]
        input_info_col_values = self.input_dict[col_values]
        input_info_row_values = self.input_dict[row_names]
        shape = slide.shapes.add_table(
            (len(input_info_row_values) + 1), len(input_info_col_names) + 1, Cm(3.62), Cm(6.2), Cm(25.5), Cm(0.5)
        )
        table = shape.table
        headers_sum = np.append([left_col_header], input_info_col_names)
        #  Fill headers
        for column_header, item in enumerate(headers_sum):
            cell = table.cell(0, column_header)
            cell.text = headers_sum[column_header]
        #  Fill left column
        for row in range(1, len(input_info_row_values) + 1):
            cell = table.cell(row, 0)
            cell.text = input_info_row_values[row - 1]
        #  Fill rest of the table
        for column in range(1, len(input_info_col_names) + 1):
            for row in range(1, (len(input_info_row_values) + 1)):
                cell = table.cell(row, column)
                value = float(input_info_col_values[column - 1][row - 1])
                cell.text = '{:,.2f}'.format(value).replace(',', '*').replace('.', ',').replace('*', '.')
        change_format_cells(table, 12)

    def make_figure(self, scenario, slide, results) -> None:
        """
        This function makes it possible to iterate over all cells in a table
        :param scenario: the selected scenario
        :param slide: selected slide where the title is needed
        :param results: type of result wanted i.e. weighted_appreciations, key outputs, appreciations
        """
        #  Fill chart with data
        chart_data = CategoryChartData()
        output_info = self.output_dict[scenario]
        chart_data.categories = output_info.keys()
        for series in list(output_info[list(output_info.keys())[0]][results].keys()):
            values = [output_info[category][results][series] for category in list(output_info.keys())]
            chart_data.add_series(series, values)
        #  implement chart on slide
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED, Cm(3.62), Cm(3.9), Cm(25.5), Cm(14.8), chart_data
        ).chart
        #  color series on chart
        colors = [
            "#D04A02",
            "#EB8C00",
            "#FFB600",
            "#295477",
            "#299D8F",
        ]
        for index, series in enumerate(chart.series):
            fill = series.format.fill
            fill.solid()
            color_used = colors[index % 5].lstrip("#")
            rgb = tuple(int(color_used[i: i + 2], 16) for i in (0, 2, 4))
            fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
        #  layout chart change
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(9)
        category_axis.tick_labels.font.bold = True
        value_axis = chart.value_axis
        value_axis.maximum_scale = 100
        value_axis.major_unit = 20
        #  titles
        chart.chart_title.text_frame.text = "Values of weighted appreciations"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        category_axis_title = chart.value_axis
        category_axis_title.axis_title.text_frame.text = "Appreciation value"
        category_axis_title.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
        #  for implementing legend
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.horz_offset = 0.0
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(8)

    def add_page_number(self, slide) -> int:
        """
        This function adds a page number to a slide and return the next page number
        :param slide: selected slide where the title is needed
        :return: the page number for the next slide
        """
        x_loc_page_number, y_loc_page_number, x_weight_page_number, y_weight_page_number = (
            Cm(31.9),
            Cm(17.7),
            Cm(0.8),
            Cm(1),
        )
        tx_box = slide.shapes.add_textbox(
            x_loc_page_number, y_loc_page_number, x_weight_page_number, y_weight_page_number
        )
        tx_box.text_frame.text = str(self.page_number)
        tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(120, 120, 120)
        self.page_number += 1
        return self.page_number

    def make_slides(self, scenario) -> Presentation():
        """
        This function present different visualisation on different slides in a PowerPoint format
        :param scenario: the desired scenario which is used in the visualisations
        :return: the created PowerPoint
        """
        prs = Presentation()
        # set width to full width.
        prs.slide_width = Inches(13)

        # Generate title slide
        slide_title = prs.slides.add_slide(prs.slide_layouts[0])
        self.make_title(slide_title, "front_page", True)
        make_subtitle(slide_title)

        # Generate strategic challenge slide
        slide_challenge = prs.slides.add_slide(prs.slide_layouts[5])
        self.make_title(slide_challenge, "strategic_challenge")
        self.make_strategic_challenge(slide_challenge)
        self.page_number = self.add_page_number(slide_challenge)

        # Create the table chart slide key outputs
        slide_key_outputs = prs.slides.add_slide(prs.slide_layouts[5])
        self.make_title(slide_key_outputs, "key_outputs")
        self.make_introduction(slide_key_outputs, "key_outputs")
        self.page_number = self.add_page_number(slide_key_outputs)
        self.make_table_2_col(slide_key_outputs, "key_outputs", "key_output_theme", "Key outputs", "Theme")

        # Create the table chart slide decision maker's options
        slide_dmo = prs.slides.add_slide(prs.slide_layouts[5])
        self.make_title(slide_dmo, "dmo")
        self.make_introduction(slide_dmo, "dmo")
        self.page_number = self.add_page_number(slide_dmo)
        self.make_table_n_col(
            slide_dmo,
            "decision_makers_options",
            "decision_makers_option_value",
            "internal_variable_inputs",
            "Internal variable input",
        )

        # Create the table chart slide the scenarios
        slide_scenarios = prs.slides.add_slide(prs.slide_layouts[5])
        self.make_title(slide_scenarios, "scenarios")
        self.make_introduction(slide_scenarios, "scenarios")
        self.page_number = self.add_page_number(slide_scenarios)
        self.make_table_n_col(
            slide_scenarios, "scenarios", "scenario_value", "external_variable_inputs", "External variable input"
        )

        # Create the table chart slide fixed inputs
        number_of_iterations = round((len(self.input_dict["fixed_inputs"]) / 10) + 0.5)
        for number_iteration in range(0, number_of_iterations):
            slide_fixed_inputs = prs.slides.add_slide(prs.slide_layouts[5])
            if number_of_iterations > 1:
                self.make_title(
                    slide_fixed_inputs,
                    "fixed_inputs",
                    pos_series=" " + str(number_iteration + 1) + "/" + str(number_of_iterations),
                )
            else:
                self.make_title(slide_fixed_inputs, "fixed_inputs")
            self.make_introduction(slide_fixed_inputs, "fixed_inputs")
            self.page_number = self.add_page_number(slide_fixed_inputs)
            self.make_table_2_col(
                slide_fixed_inputs,
                "fixed_inputs",
                "fixed_input_value",
                "Fixed Inputs",
                "Value",
                (number_iteration * 10),
                (number_iteration * 10) + 10,
            )

        # Generate dependency graph slide
        rand_dir = random()
        os.mkdir(str(rand_dir))
        for key_output in self.output_dict[scenario][next(iter(self.output_dict[scenario]))]["key_outputs"].items():
            slide_dep_graph = prs.slides.add_slide(prs.slide_layouts[5])
            self.make_title(slide_dep_graph, "dependency_graph", key_output=" for key output: " + key_output[0])
            self.make_introduction(slide_dep_graph, "dependency_graph")
            # slide_dep_graph.shapes.add_picture("dep-graph.png", Cm(3.62), Cm(6.2), Cm(25.5), Cm(12.5))
            self.page_number = self.add_page_number(slide_dep_graph)
        os.rmdir(str(rand_dir))

        # Generate weighted appreciations slide
        slide_weighted_app = prs.slides.add_slide(prs.slide_layouts[5])
        self.make_title(slide_weighted_app, "weighted_graph", scenario=scenario)
        self.add_page_number(slide_weighted_app)
        self.make_figure(scenario, slide_weighted_app, "weighted_appreciations")

        return prs

    def create_report(self, scenario, path) -> str:
        """
        This function saves the created report out of the function make_slides at the desired location
        :param scenario: the desired scenario which is used in the visualisations
        :param path: the desired location where the report is saved
        :return: the text where the PowerPoint is saved
        """
        prs = self.make_slides(scenario)
        date_year = str(datetime.now().strftime("%Y-%m-%d"))
        date_hour = str(datetime.now().strftime("%Hh%Mm%Ss"))
        filename = "tRBS_" + self.name + f'_{date_year + "_" + date_hour}.pptx'
        prs.save(str(path) + "/" + filename)
        text_finished = "The powerpoint is generated and located at " + str(path) + "/" + filename
        return text_finished
